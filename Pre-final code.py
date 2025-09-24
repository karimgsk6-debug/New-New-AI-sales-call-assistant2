# app.py
import os
import io
import base64
import requests
from pathlib import Path
from datetime import datetime
import asyncio

import streamlit as st
from PIL import Image
from docx import Document
import pdfplumber
from pptx import Presentation

from groq import Groq
import edge_tts

# ----------------------------
# App configuration
# ----------------------------
st.set_page_config(page_title="AI Sales Call Assistant", layout="wide", initial_sidebar_state="expanded")

# ----------------------------
# Constants
# ----------------------------
BACKGROUND_URL = "https://image.shutterstock.com/image-photo/young-arab-girl-using-ipad-260nw-2616487693.jpg"
GROQ_API_KEY = os.getenv("GROQ_API_KEY", "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")

# ----------------------------
# Helper functions
# ----------------------------
def extract_text_from_docx(file):
    doc = Document(file)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pdf(file):
    text = ""
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

async def generate_tts_edge_async(text, voice="en-US-JennyNeural", filename=None):
    if filename is None:
        filename = f"ai_tts_{datetime.now().strftime('%Y%m%d_%H%M%S%f')}.mp3"
    communicate = edge_tts.Communicate(text, voice=voice)
    await communicate.save(filename)
    return filename

def generate_tts_edge(text, voice="en-US-JennyNeural", filename=None):
    return asyncio.run(generate_tts_edge_async(text, voice=voice, filename=filename))

def safe_groq_client():
    if GROQ_API_KEY:
        try:
            return Groq(api_key=GROQ_API_KEY)
        except Exception as e:
            st.warning(f"Could not initialize Groq client: {e}")
            return None
    return None

def ask_ai_via_groq(prompt, client=None, fallback_message="⚠️ Groq API not configured or request failed."):
    if client is None:
        client = safe_groq_client()
    if client is None:
        return fallback_message
    try:
        resp = client.chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=[
                {"role": "system", "content": "You are a helpful AI medical sales assistant. Structure responses according to the pharma sales call flow. Use APACT only when handling objections and highlight each step. Reference uploaded docs if available."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.6,
            max_tokens=1000,
        )
        return resp.choices[0].message.content
    except Exception as e:
        return f"{fallback_message} Error: {e}"

# ----------------------------
# HEADER: Page title + disclaimer
# ----------------------------
st.markdown(f"""
<div style='text-align:center; padding:15px; background:linear-gradient(90deg,#ff8c00,#ffb347); 
            color:white; border-radius:12px; margin-bottom:10px;'>
    <h2 style='margin:0;'>💡 AI Sales Call Assistant</h2>
    <p style='margin:0;'>Powered by AI to equip sales reps for smarter HCP conversations</p>
</div>

<div style='padding:10px; background:#ffffff; border:1px solid #ddd; border-radius:10px; margin-bottom:20px; font-size:13px; color:#000000;'>
    ⚠️ <b>Disclaimer:</b> This AI tool is to equip sales reps and is not a substitute for official product info or medical advice.
</div>
""", unsafe_allow_html=True)

# ----------------------------
# Sidebar: filters
# ----------------------------
st.sidebar.header("⚙️ Settings & Filters")

# Theme toggle
theme_choice = st.sidebar.radio("Theme", options=["Dark Mode", "Light Mode"], index=1)

# Brand selection
st.sidebar.subheader("Brand & Segmentation")
gsk_brands = {
    "Shingrix": "https://example.com/shingrix-leaflet",
    "Trelegy": "https://example.com/trelegy-leaflet",
    "Zejula": "https://example.com/zejula-leaflet",
}
brand = st.sidebar.selectbox("💊 Select Brand", options=list(gsk_brands.keys()))

# HCP Segmentation
st.sidebar.subheader("HCP Segmentation")
hcp_segments = [
    "Unaware: HCP has no knowledge of disease/product",
    "Apathetic: Aware but not engaged",
    "Neutral: Aware, waiting for more evidence",
    "Supportive: Positive but not proactive",
    "Advocate: Actively prescribes and promotes",
]
segment = st.sidebar.selectbox("👥 Segment", hcp_segments)

# Barriers
st.sidebar.subheader("Doctor Barriers")
doctor_barriers = [
    "HCP does not consider HZ as risk",
    "No time to discuss preventive measures",
    "Cost considerations",
    "Not convinced HZ Vx effective",
    "Accessibility issues",
    "Regulatory concerns",
    "Patient hesitancy",
]
barrier = st.sidebar.multiselect("🚧 Select Barriers", options=doctor_barriers, default=[])

# Attributes
st.sidebar.subheader("Doctor / HCP Attributes")
specialty = st.sidebar.selectbox("🩺 Specialty", options=["GP","Cardiologist","Dermatologist","Endocrinologist","Pulmonologist","Immunologist"])
persona = st.sidebar.selectbox("🧑‍⚕️ Persona", options=["Uncommitted Vaccinator","Reluctant Efficiency","Patient Influenced","Committed Vaccinator"])

# Updated tone & mindset filters
response_tone = st.sidebar.selectbox("🎤 Response Tone", options=["Formal","Empathetic","Confident","Concise","Persuasive"])
hcp_mindset = st.sidebar.selectbox("💡 HCP Mindset", options=["Analytical","Practical","Risk-Averse","Innovative","Skeptical","Patient-Centered"])

# Call stage
st.sidebar.markdown("---")
call_stage = st.sidebar.selectbox("📞 Call Stage", options=[
    "Prepare the Call","Engage","Create Opportunities","Influence","Impact GSO (Good Sell Outcome)",
    "Closing with Commitment","Post-Call Analysis"
])

# ----------------------------
# Adaptive colors
# ----------------------------
is_light = theme_choice == "Light Mode"
font_color = "black" if is_light else "white"
bubble_user_bg = "rgba(255,255,255,0.85)" if is_light else "rgba(255,255,255,0.14)"
bubble_ai_bg = "rgba(255,255,255,0.65)" if is_light else "rgba(0,0,0,0.35)"

# ----------------------------
# Full-page blurred background
# ----------------------------
st.markdown(f"""
<style>
[data-testid="stAppViewContainer"] {{
    position: relative;
    background-image: url("{BACKGROUND_URL}");
    background-repeat: no-repeat;
    background-position: center top;
    background-size: cover;
    filter: blur(3px);
    height: 100vh;
    width: 100%;
    z-index: -1;
    position: fixed;
}}

/* Chat container and inputs overlay */
.chat-container, .stTextInput, .stTextArea {{
    background: rgba(0,0,0,0.5) !important;
    color: {font_color} !important;
    border-radius: 12px;
    padding: 10px;
}}
.user-bubble {{
    text-align:right;
    background:{bubble_user_bg};
    color:{font_color};
    padding:10px;
    border-radius:15px 15px 0 15px;
    margin:6px;
    display:inline-block;
    max-width:80%;
}}
.ai-bubble {{
    text-align:left;
    background:{bubble_ai_bg};
    color:{font_color};
    padding:10px;
    border-radius:15px 15px 15px 0;
    margin:6px;
    display:inline-block;
    max-width:80%;
}}
.apact-step {{
    background:#ffd700; color:#000; font-weight:bold; padding:2px 6px; border-radius:4px;
}}
</style>
""", unsafe_allow_html=True)

# ----------------------------
# Chat history
# ----------------------------
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "uploaded_docs" not in st.session_state:
    st.session_state.uploaded_docs = ""

chat_placeholder = st.empty()
def display_chat():
    chat_html = "<div class='chat-container'>"
    for msg in st.session_state.chat_history:
        content = msg["content"].replace("\n","<br>")
        for step in ["Acknowledge","Probing","Action","Confirm","Transition"]:
            content = content.replace(step, f"<span class='apact-step'>{step}</span>")
        if msg["role"] == "user":
            chat_html += f"<div class='user-bubble'>{content}</div>"
        else:
            chat_html += f"<div class='ai-bubble'>{content}</div>"
    chat_html += "</div>"
    chat_placeholder.markdown(chat_html, unsafe_allow_html=True)
display_chat()

# ----------------------------
# Input
# ----------------------------
user_input = st.text_input("💬 Type your message:", "")

if st.button("📩 Send") and user_input:
    st.session_state.chat_history.append({"role":"user","content":user_input})
    display_chat()

    # AI response
    prompt = f"""
Stage: {call_stage}
Segment: {segment}
Barriers: {', '.join(barrier) if barrier else 'None'}
Brand: {brand}
Specialty: {specialty}
Persona: {persona}
HCP Mindset: {hcp_mindset}
Tone: {response_tone}
Docs: {st.session_state.uploaded_docs[:1000]}
Input: {user_input}
"""
    groq_client = safe_groq_client()
    ai_text = ask_ai_via_groq(prompt, groq_client)
    st.session_state.chat_history.append({"role":"ai","content":ai_text})
    display_chat()
